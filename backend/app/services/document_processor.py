import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
import openpyxl
import markdown
from pathlib import Path
from typing import Dict, List
import logging
import zipfile
import tarfile
import json
import yaml
import configparser

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Extraction de texte multi-format (30+ formats supportés)"""
    
    SUPPORTED_FORMATS = {
        '.pdf': 'PDF Document',
        '.docx': 'Word Document',
        '.doc': 'Word Document (Legacy)',
        '.pptx': 'PowerPoint Presentation',
        '.ppt': 'PowerPoint (Legacy)',
        '.txt': 'Text File',
        '.md': 'Markdown File',
        '.markdown': 'Markdown File',
        '.html': 'HTML File',
        '.htm': 'HTML File',
        '.xlsx': 'Excel Spreadsheet',
        '.xls': 'Excel (Legacy)',
        '.csv': 'CSV File',
        '.rtf': 'Rich Text Format',
        '.odt': 'OpenDocument Text',
        '.ods': 'OpenDocument Spreadsheet',
        '.odp': 'OpenDocument Presentation',
        '.tex': 'LaTeX Document',
        '.epub': 'EPUB eBook',
        '.xml': 'XML File',
        
        # Code - Python & Web
        '.py': 'Python Script',
        '.js': 'JavaScript File',
        '.jsx': 'React JSX File',
        '.ts': 'TypeScript File',
        '.tsx': 'React TSX File',
        '.css': 'CSS Stylesheet',
        
        # Code - Backend
        '.java': 'Java Source File',
        '.cpp': 'C++ Source File',
        '.c': 'C Source File',
        '.cs': 'C# Source File',
        '.go': 'Go Source File',
        '.rs': 'Rust Source File',
        '.php': 'PHP Script',
        '.rb': 'Ruby Script',
        '.swift': 'Swift Source File',
        '.kt': 'Kotlin Source File',
        '.scala': 'Scala Source File',
        '.r': 'R Script',
        '.groovy': 'Groovy Script',
        
        # Scripts & Config
        '.sh': 'Shell Script',
        '.bash': 'Bash Script',
        '.sql': 'SQL Script',
        '.json': 'JSON File',
        '.yaml': 'YAML Config',
        '.yml': 'YAML Config',
        '.toml': 'TOML Config',
        '.ini': 'INI Config',
        '.env': 'Environment File',
        '.jenkinsfile': 'Jenkins Pipeline',
        
        # Archives
        '.zip': 'ZIP Archive',
        '.tar': 'TAR Archive',
        '.gz': 'GZIP Archive',
        '.tar.gz': 'TAR.GZ Archive',
    }
    
    @staticmethod
    def is_supported(filename: str) -> bool:
        """Vérifie si le format est supporté"""
        path = Path(filename)
        
        # Cas spécial pour .tar.gz
        if filename.endswith('.tar.gz'):
            return True
        
        suffix = path.suffix.lower()
        return suffix in DocumentProcessor.SUPPORTED_FORMATS
    
    @staticmethod
    def extract_text(file_path: Path) -> str:
        """Route vers le bon extracteur"""
        # Cas spécial pour .tar.gz
        if str(file_path).endswith('.tar.gz'):
            logger.info(f"📦 Extracting TAR.GZ archive: {file_path.name}")
            return DocumentProcessor._extract_tar_gz(file_path)
        
        suffix = file_path.suffix.lower()
        
        # Map des extracteurs
        extractors = {
            # Documents
            '.pdf': DocumentProcessor._extract_pdf,
            '.docx': DocumentProcessor._extract_docx,
            '.doc': DocumentProcessor._extract_code, 
            '.pptx': DocumentProcessor._extract_pptx,
            '.ppt': DocumentProcessor._extract_code,
            '.txt': DocumentProcessor._extract_txt,
            '.md': DocumentProcessor._extract_markdown,
            '.markdown': DocumentProcessor._extract_markdown,
            '.html': DocumentProcessor._extract_html,
            '.htm': DocumentProcessor._extract_html,
            '.xlsx': DocumentProcessor._extract_excel,
            '.xls': DocumentProcessor._extract_code,
            '.csv': DocumentProcessor._extract_csv,
            '.rtf': DocumentProcessor._extract_code,
            '.odt': DocumentProcessor._extract_code,
            '.ods': DocumentProcessor._extract_code,
            '.odp': DocumentProcessor._extract_code,
            '.tex': DocumentProcessor._extract_code,
            '.epub': DocumentProcessor._extract_code,
            '.xml': DocumentProcessor._extract_code,
            
            # Code (tous gérés comme du texte brut)
            '.py': DocumentProcessor._extract_code,
            '.js': DocumentProcessor._extract_code,
            '.jsx': DocumentProcessor._extract_code,
            '.ts': DocumentProcessor._extract_code,
            '.tsx': DocumentProcessor._extract_code,
            '.css': DocumentProcessor._extract_code,
            '.java': DocumentProcessor._extract_code,
            '.cpp': DocumentProcessor._extract_code,
            '.c': DocumentProcessor._extract_code,
            '.cs': DocumentProcessor._extract_code,
            '.go': DocumentProcessor._extract_code,
            '.rs': DocumentProcessor._extract_code,
            '.php': DocumentProcessor._extract_code,
            '.rb': DocumentProcessor._extract_code,
            '.swift': DocumentProcessor._extract_code,
            '.kt': DocumentProcessor._extract_code,
            '.scala': DocumentProcessor._extract_code,
            '.r': DocumentProcessor._extract_code,
            '.groovy': DocumentProcessor._extract_code,
            '.sh': DocumentProcessor._extract_code,
            '.bash': DocumentProcessor._extract_code,
            '.sql': DocumentProcessor._extract_code,
            
            # Config
            '.json': DocumentProcessor._extract_json,
            '.yaml': DocumentProcessor._extract_yaml,
            '.yml': DocumentProcessor._extract_yaml,
            '.toml': DocumentProcessor._extract_toml,
            '.ini': DocumentProcessor._extract_ini,
            '.env': DocumentProcessor._extract_code,
            '.jenkinsfile': DocumentProcessor._extract_code,
            
            # Archives
            '.zip': DocumentProcessor._extract_zip,
            '.tar': DocumentProcessor._extract_tar,
            '.gz': DocumentProcessor._extract_gz,
        }
        
        extractor = extractors.get(suffix)
        if not extractor:
            raise ValueError(f"Format non supporté : {suffix}")
        
        logger.info(f"📄 Extracting {DocumentProcessor.SUPPORTED_FORMATS.get(suffix, 'file')}: {file_path.name}")
        return extractor(file_path)
    
    # ================== DOCUMENTS ==================
    
    @staticmethod
    def _extract_pdf(path: Path) -> str:
        """Extraction PDF avec pypdf"""
        text = []
        try:
            with open(path, 'rb') as f:
                pdf = pypdf.PdfReader(f)
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text.append(f"--- Page {i+1} ---\n{page_text}")
            return "\n\n".join(text)
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_docx(path: Path) -> str:
        """Extraction DOCX"""
        try:
            doc = DocxDocument(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Extraction des tableaux
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(" | ".join(row_data))
                tables.append("\n".join(table_data))
            
            content = "\n\n".join(paragraphs)
            if tables:
                content += "\n\n--- Tables ---\n\n" + "\n\n".join(tables)
            
            return content
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_pptx(path: Path) -> str:
        """Extraction PowerPoint"""
        try:
            prs = Presentation(path)
            text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {i+1} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                text.append("\n".join(slide_text))
            
            return "\n\n".join(text)
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_txt(path: Path) -> str:
        """Extraction fichier texte"""
        try:
            return path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            # Fallback latin-1
            return path.read_text(encoding='latin-1')
    
    @staticmethod
    def _extract_markdown(path: Path) -> str:
        """Extraction Markdown (garde le markdown brut pour le code)"""
        try:
            return path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            return path.read_text(encoding='latin-1')
    
    @staticmethod
    def _extract_html(path: Path) -> str:
        """Extraction HTML"""
        try:
            html = path.read_text(encoding='utf-8')
            soup = BeautifulSoup(html, 'html.parser')
            
            # Supprimer scripts et styles
            for script in soup(["script", "style"]):
                script.decompose()
            
            return soup.get_text(separator="\n", strip=True)
        except Exception as e:
            logger.error(f"HTML extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_excel(path: Path) -> str:
        """Extraction Excel"""
        try:
            wb = openpyxl.load_workbook(path, data_only=True)
            text = []
            
            for sheet in wb.worksheets:
                text.append(f"--- Sheet: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text.append(row_text)
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Excel extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_csv(path: Path) -> str:
        """Extraction CSV"""
        try:
            import csv
            text = []
            with open(path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row in reader:
                    text.append(" | ".join(row))
            return "\n".join(text)
        except Exception as e:
            logger.error(f"CSV extraction error: {e}")
            raise
    
    # ================== CODE ==================
    
    @staticmethod
    def _extract_code(path: Path) -> str:
        """Extraction de fichiers code (préserve la syntaxe)"""
        try:
            content = path.read_text(encoding='utf-8')
            
            # Ajouter un header avec le nom du fichier et le langage
            file_type = DocumentProcessor.SUPPORTED_FORMATS.get(path.suffix.lower(), 'Code File')
            header = f"=== {path.name} ({file_type}) ===\n\n"
            
            return header + content
        except UnicodeDecodeError:
            try:
                content = path.read_text(encoding='latin-1')
                file_type = DocumentProcessor.SUPPORTED_FORMATS.get(path.suffix.lower(), 'Code File')
                header = f"=== {path.name} ({file_type}) ===\n\n"
                return header + content
            except Exception as e:
                logger.error(f"Code extraction error: {e}")
                raise
    
    # ================== CONFIG ==================
    
    @staticmethod
    def _extract_json(path: Path) -> str:
        """Extraction JSON formaté"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Reformater proprement pour la lisibilité
            formatted = json.dumps(data, indent=2, ensure_ascii=False)
            return f"=== {path.name} (JSON Config) ===\n\n{formatted}"
        except Exception as e:
            logger.error(f"JSON extraction error: {e}")
            # Fallback: lire comme texte brut
            return DocumentProcessor._extract_code(path)
    
    @staticmethod
    def _extract_yaml(path: Path) -> str:
        """Extraction YAML"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
            
            # Convertir en texte lisible
            formatted = yaml.dump(data, default_flow_style=False, allow_unicode=True)
            return f"=== {path.name} (YAML Config) ===\n\n{formatted}"
        except Exception as e:
            logger.warning(f"YAML parsing failed, reading as text: {e}")
            # Fallback: lire comme texte brut
            return DocumentProcessor._extract_code(path)
    
    @staticmethod
    def _extract_toml(path: Path) -> str:
        """Extraction TOML"""
        try:
            import tomli  # Python 3.11+ a tomllib intégré
            
            with open(path, 'rb') as f:
                data = tomli.load(f)
            
            # Convertir en JSON pour affichage
            formatted = json.dumps(data, indent=2, ensure_ascii=False)
            return f"=== {path.name} (TOML Config) ===\n\n{formatted}"
        except ImportError:
            logger.warning("tomli not installed, reading TOML as text")
            return DocumentProcessor._extract_code(path)
        except Exception as e:
            logger.warning(f"TOML parsing failed: {e}")
            return DocumentProcessor._extract_code(path)
    
    @staticmethod
    def _extract_ini(path: Path) -> str:
        """Extraction INI"""
        try:
            config = configparser.ConfigParser()
            config.read(path, encoding='utf-8')
            
            text = [f"=== {path.name} (INI Config) ===\n"]
            for section in config.sections():
                text.append(f"\n[{section}]")
                for key, value in config.items(section):
                    text.append(f"{key} = {value}")
            
            return "\n".join(text)
        except Exception as e:
            logger.warning(f"INI parsing failed: {e}")
            return DocumentProcessor._extract_code(path)
    
    # ================== ARCHIVES ==================
    
    @staticmethod
    def _extract_zip(path: Path) -> str:
        """Extraction ZIP (liste des fichiers + extraction texte)"""
        try:
            text = [f"=== ZIP Archive: {path.name} ===\n"]
            
            with zipfile.ZipFile(path, 'r') as zip_file:
                file_list = zip_file.namelist()
                text.append(f"📦 Contains {len(file_list)} files:\n")
                
                for file_name in file_list[:50]:  # Limite à 50 fichiers
                    text.append(f"  - {file_name}")
                
                if len(file_list) > 50:
                    text.append(f"  ... and {len(file_list) - 50} more files")
                
                # Extraire les fichiers texte (max 10)
                text.append("\n📄 Text content from archive:\n")
                extracted_count = 0
                
                for file_name in file_list:
                    if extracted_count >= 10:
                        break
                    
                    # Vérifier si c'est un fichier supporté
                    if DocumentProcessor.is_supported(file_name):
                        try:
                            file_data = zip_file.read(file_name)
                            
                            # Essayer de décoder comme texte
                            try:
                                file_text = file_data.decode('utf-8')
                                text.append(f"\n--- {file_name} ---\n{file_text[:1000]}")  # Max 1000 chars/file
                                extracted_count += 1
                            except UnicodeDecodeError:
                                pass
                        except Exception as e:
                            logger.warning(f"Could not extract {file_name}: {e}")
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"ZIP extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_tar(path: Path) -> str:
        """Extraction TAR"""
        try:
            text = [f"=== TAR Archive: {path.name} ===\n"]
            
            with tarfile.open(path, 'r') as tar:
                members = tar.getmembers()
                text.append(f"📦 Contains {len(members)} files:\n")
                
                for member in members[:50]:
                    text.append(f"  - {member.name}")
                
                if len(members) > 50:
                    text.append(f"  ... and {len(members) - 50} more files")
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"TAR extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_tar_gz(path: Path) -> str:
        """Extraction TAR.GZ"""
        try:
            text = [f"=== TAR.GZ Archive: {path.name} ===\n"]
            
            with tarfile.open(path, 'r:gz') as tar:
                members = tar.getmembers()
                text.append(f"📦 Contains {len(members)} files:\n")
                
                for member in members[:50]:
                    text.append(f"  - {member.name}")
                
                if len(members) > 50:
                    text.append(f"  ... and {len(members) - 50} more files")
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"TAR.GZ extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_gz(path: Path) -> str:
        """Extraction GZIP simple"""
        try:
            import gzip
            
            with gzip.open(path, 'rt', encoding='utf-8') as f:
                content = f.read()
            
            return f"=== GZIP File: {path.name} ===\n\n{content}"
        except Exception as e:
            logger.error(f"GZIP extraction error: {e}")
            raise
    
    # ================== UTILITIES ==================
    
    @classmethod
    def get_format_category(cls, filename: str) -> str:
        """Retourne la catégorie du format"""
        suffix = Path(filename).suffix.lower()
        
        if suffix in ['.py', '.js', '.jsx', '.ts', '.tsx', '.java', '.cpp', '.c', '.cs', 
                      '.go', '.rs', '.php', '.rb', '.swift', '.kt', '.scala', '.r']:
            return "code"
        
        if suffix in ['.json', '.yaml', '.yml', '.toml', '.ini', '.env']:
            return "config"
        
        if suffix in ['.zip', '.tar', '.gz'] or filename.endswith('.tar.gz'):
            return "archive"
        
        if suffix in ['.md', '.txt', '.html', '.htm']:
            return "document"
        
        if suffix in ['.pdf', '.docx', '.pptx', '.xlsx', '.csv']:
            return "office"
        
        return "unknown"
    
    @classmethod
    def get_recommended_chunk_profile(cls, filename: str, file_size_bytes: int) -> str:
        """Recommande un profil de chunking selon le type de fichier"""
        category = cls.get_format_category(filename)
        suffix = Path(filename).suffix.lower()
        
        # Code → gros chunks pour garder la structure
        if category == "code":
            if file_size_bytes > 100_000:  # > 100KB
                return "project_full"
            return "code_analysis"
        
        # Config → chunks moyens
        if category == "config":
            return "standard"
        
        # Archives → chunks moyens (juste la liste)
        if category == "archive":
            return "standard"
        
        # Documentation → adaptatif selon taille
        if category == "document":
            if file_size_bytes > 200_000:  # > 200KB
                return "project_full"
            return "documentation"
        
        # Office → standard
        return "standard"